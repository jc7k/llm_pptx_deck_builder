#!/usr/bin/env python3
"""Test bullet point formatting in presentations."""

import os
import sys
from datetime import datetime
from pptx import Presentation

# Add src to path
sys.path.insert(0, 'src')

def test_bullet_formatting():
    """Test bullet point formatting without double bullets."""
    
    # Create a test presentation
    prs = Presentation()
    
    # Create title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Test Presentation"
    subtitle = title_slide.placeholders[1]
    subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    
    # Create agenda slide
    agenda_slide = prs.slides.add_slide(prs.slide_layouts[1])
    agenda_slide.shapes.title.text = "Agenda"
    
    agenda_content = agenda_slide.placeholders[1]
    agenda_items = ["Introduction", "Key Concepts", "Current Trends"]
    # Test without manual bullets - let the slide layout handle it
    agenda_content.text = "\n".join(agenda_items)
    
    # Create content slide
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    content_slide.shapes.title.text = "Introduction"
    
    content = content_slide.placeholders[1]
    bullets = [
        "AI is projected to create 170 million new jobs by 2030",
        "75% of knowledge workers are adopting AI tools",
        "AI-skilled workers earn a 56% wage premium"
    ]
    # Test without manual bullets
    content.text = "\n".join(bullets)
    
    # Save test presentation
    output_path = "output/test_bullets.pptx"
    os.makedirs("output", exist_ok=True)
    prs.save(output_path)
    
    print(f"✅ Test presentation saved to: {output_path}")
    print("Please check if bullet points appear correctly (single bullets, not double)")
    
    return output_path

if __name__ == "__main__":
    test_bullet_formatting()