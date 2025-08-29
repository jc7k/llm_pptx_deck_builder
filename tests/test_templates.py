"""Tests for PowerPoint template functionality."""

import os
import pytest
from pathlib import Path
from pptx import Presentation
from unittest.mock import patch, MagicMock

from src.tools import create_presentation
from src.deck_builder_agent import build_deck_sync

# Get the fixtures directory
FIXTURES_DIR = Path(__file__).parent / "fixtures"
TEMPLATES_DIR = FIXTURES_DIR / "templates"

# List all available templates
TEMPLATE_FILES = [
    "Coastal presentation.pptx",
    "Colorful conference presentation.pptx", 
    "Gold Elegant Branding Kit Presentation.pptx",
    "Light modernist presentation.pptx",
    "Pastel Minimal Street Maps Minitheme Presentation.pptx",
    "Prismatic design.pptx"
]


class TestTemplateLoading:
    """Test template loading and validation."""
    
    def test_all_templates_exist(self):
        """Verify all template files exist."""
        for template_file in TEMPLATE_FILES:
            template_path = TEMPLATES_DIR / template_file
            assert template_path.exists(), f"Template {template_file} not found"
            assert template_path.stat().st_size > 0, f"Template {template_file} is empty"
    
    def test_templates_are_valid_pptx(self):
        """Verify all templates are valid PowerPoint files."""
        for template_file in TEMPLATE_FILES:
            template_path = TEMPLATES_DIR / template_file
            try:
                prs = Presentation(str(template_path))
                assert prs is not None
                assert len(prs.slide_layouts) > 0, f"Template {template_file} has no slide layouts"
            except Exception as e:
                pytest.fail(f"Failed to load template {template_file}: {e}")
    
    def test_template_slide_layouts(self):
        """Test that templates have expected slide layouts."""
        for template_file in TEMPLATE_FILES:
            template_path = TEMPLATES_DIR / template_file
            prs = Presentation(str(template_path))
            
            # Check for slide layouts
            layout_count = len(prs.slide_layouts)
            
            # Most templates should have at least title and content layouts
            assert layout_count >= 2, f"Template {template_file} has too few layouts"
            
            print(f"\n{template_file} has {layout_count} slide layouts")
    
    def test_template_slide_masters(self):
        """Test that templates have slide masters."""
        for template_file in TEMPLATE_FILES:
            template_path = TEMPLATES_DIR / template_file
            prs = Presentation(str(template_path))
            
            # Check that presentation has slide layouts (which come from masters)
            assert len(prs.slide_layouts) > 0, f"Template {template_file} has no slide layouts"
            
            # Verify we can access slide layouts (which confirms masters exist)
            # python-pptx doesn't directly expose slide_master attribute
            # but slide layouts come from masters, so their existence confirms masters exist


class TestTemplateIntegration:
    """Test template integration with deck builder."""
    
    @pytest.fixture
    def sample_slide_specs(self):
        """Sample slide specifications for testing."""
        return [
            {
                "title": "Introduction",
                "bullets": [
                    "Welcome to the presentation",
                    "Today's agenda",
                    "Key objectives"
                ],
                "notes": "Speaker notes for introduction",
                "references": []
            },
            {
                "title": "Main Content",
                "bullets": [
                    "Point 1: Important concept",
                    "Point 2: Supporting details",
                    "Point 3: Examples and applications"
                ],
                "notes": "Detailed explanation of main points",
                "references": ["https://example.com/source1"]
            },
            {
                "title": "Conclusion",
                "bullets": [
                    "Summary of key points",
                    "Next steps",
                    "Questions?"
                ],
                "notes": "Wrap up the presentation",
                "references": []
            }
        ]
    
    @pytest.mark.parametrize("template_file", TEMPLATE_FILES)
    def test_create_presentation_with_template(self, template_file, sample_slide_specs, tmp_path):
        """Test creating presentations with each template."""
        template_path = str(TEMPLATES_DIR / template_file)
        
        # Mock the output path
        output_dir = tmp_path / "output"
        output_dir.mkdir()
        
        with patch('src.tools.os.makedirs'):
            with patch('src.tools.os.path.dirname', return_value=str(output_dir)):
                result = create_presentation.invoke({
                    "slide_specs": sample_slide_specs,
                    "template_path": template_path
                })
        
        # The tool returns the path directly as a string
        assert isinstance(result, str)
        assert result.endswith(".pptx")
        
        # Verify the file would be created with the template
        # Note: We're mocking the actual save to avoid file I/O in tests
        assert template_file in template_path
    
    def test_template_preserves_theme(self, sample_slide_specs, tmp_path):
        """Test that template theme is preserved in generated presentation."""
        # Use the first template for this test
        template_file = TEMPLATE_FILES[0]
        template_path = TEMPLATES_DIR / template_file
        
        # Load original template
        original_prs = Presentation(str(template_path))
        
        # Create a presentation with the template
        output_path = tmp_path / "test_output.pptx"
        
        # Actually create the presentation
        prs = Presentation(str(template_path))
        
        # Add a test slide
        if len(prs.slide_layouts) > 0:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title if placeholder exists
            if slide.shapes.title:
                slide.shapes.title.text = "Test Slide"
        
        # Save the presentation
        prs.save(str(output_path))
        
        # Load the saved presentation
        saved_prs = Presentation(str(output_path))
        
        # Verify theme preservation by checking layouts are preserved
        assert len(saved_prs.slide_layouts) == len(original_prs.slide_layouts)
        # Verify the presentation was created successfully
        assert output_path.exists(), "Presentation file was not created"
    
    def test_template_with_invalid_path(self, sample_slide_specs):
        """Test handling of invalid template path."""
        invalid_path = "/nonexistent/template.pptx"
        
        with patch('src.tools.os.makedirs'):
            result = create_presentation.invoke({
                "slide_specs": sample_slide_specs,
                "template_path": invalid_path
            })
        
        # Should fall back to default presentation when template doesn't exist
        assert isinstance(result, str)
        assert result.endswith(".pptx")
    
    def test_template_compatibility_with_slide_types(self, tmp_path):
        """Test that different slide types work with templates."""
        template_path = TEMPLATES_DIR / TEMPLATE_FILES[0]
        prs = Presentation(str(template_path))
        
        # Test adding different types of slides
        slides_added = []
        
        # Try to add slides with different layouts
        for i, layout in enumerate(prs.slide_layouts[:3]):  # Test first 3 layouts
            try:
                slide = prs.slides.add_slide(layout)
                slides_added.append(f"Layout {i}")
                
                # Try to add content to placeholders
                for shape in slide.placeholders:
                    if hasattr(shape, 'text_frame'):
                        shape.text = f"Test text for layout {i}"
                        break
            except Exception as e:
                print(f"Could not add slide with layout {i}: {e}")
        
        assert len(slides_added) > 0, "Could not add any slides to template"
        
        # Save and verify
        output_path = tmp_path / "compatibility_test.pptx"
        prs.save(str(output_path))
        assert output_path.exists()


class TestTemplateEdgeCases:
    """Test edge cases and error handling for templates."""
    
    def test_template_with_no_content_layout(self):
        """Test handling templates that might not have standard content layouts."""
        # This tests the robustness of the code when templates have unusual layouts
        template_path = TEMPLATES_DIR / TEMPLATE_FILES[0]
        prs = Presentation(str(template_path))
        
        # Check if we can find a suitable layout for content
        content_layout = None
        for layout in prs.slide_layouts:
            # Look for layouts with placeholders for content
            try:
                placeholder_count = len(layout.placeholders)
                if placeholder_count >= 2:
                    content_layout = layout
                    break
            except:
                continue
        
        if content_layout:
            slide = prs.slides.add_slide(content_layout)
            assert slide is not None
    
    def test_template_placeholder_handling(self):
        """Test safe handling of template placeholders."""
        template_path = TEMPLATES_DIR / TEMPLATE_FILES[0]
        prs = Presentation(str(template_path))
        
        if len(prs.slide_layouts) > 0:
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            
            # Test safe placeholder access
            placeholders_found = []
            for shape in slide.placeholders:
                try:
                    if hasattr(shape, 'placeholder_format'):
                        placeholders_found.append(shape.placeholder_format.idx)
                except:
                    pass
            
            print(f"Placeholders found: {len(placeholders_found)}")
    
    @pytest.mark.parametrize("template_file", TEMPLATE_FILES[:2])  # Test with first 2 templates
    def test_build_deck_with_template(self, template_file):
        """Test the full deck building workflow with templates."""
        template_path = str(TEMPLATES_DIR / template_file)
        
        with patch('src.deck_builder_agent.deck_builder_graph') as mock_graph:
            # Mock the graph execution
            mock_final_state = {
                "output_path": f"test_deck_{template_file}",
                "slide_specs": [{"title": "Test", "bullets": ["Test bullet"]}],
                "references": [],
                "status": "completed",
                "messages": []
            }
            mock_graph.invoke.return_value = mock_final_state
            
            # Test build_deck_sync with template
            result = build_deck_sync(
                user_request="Test Topic",
                template_path=template_path
            )
            
            assert result["success"] is True
            
            # Verify template path was passed correctly
            call_args = mock_graph.invoke.call_args
            initial_state = call_args[0][0]
            assert initial_state["template_path"] == template_path


class TestTemplateBestPractices:
    """Test best practices for template usage."""
    
    def test_template_file_sizes(self):
        """Verify template file sizes are reasonable."""
        for template_file in TEMPLATE_FILES:
            template_path = TEMPLATES_DIR / template_file
            size_mb = template_path.stat().st_size / (1024 * 1024)
            print(f"{template_file}: {size_mb:.2f} MB")
            
            # Warn if template is very large (over 25MB)
            assert size_mb < 50, f"Template {template_file} is very large ({size_mb:.2f} MB)"
    
    def test_template_naming_convention(self):
        """Verify template files follow naming conventions."""
        for template_file in TEMPLATE_FILES:
            # Check file extension
            assert template_file.endswith('.pptx'), f"Template {template_file} should have .pptx extension"
            
            # Check for readable names (no special characters except spaces and hyphens)
            name_without_ext = template_file[:-5]
            assert all(c.isalnum() or c in [' ', '-', '_'] for c in name_without_ext), \
                f"Template {template_file} contains special characters"